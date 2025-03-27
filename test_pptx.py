class myClass:
    def __init__(self, p1:int, p2:int, p3:int, p4:int):
        self.A = p1
        self.B = p2
        self.C = p3
        self.D = p4
    def __str__(self):
        return str(self.__class__)+'\t'+str(self.A)+'\t'+str(self.B)+'\t'+str(self.C)+'\t'+str(self.D)

def copy_ps(src:myClass, dest:myClass):
    dest.A = src.A
    dest.B = src.B
    dest.C = src.C
    dest.D = src.D

X = myClass(1,2,3,4)
Y = myClass(7,8,9,10)
print(f'point 1, X={X}, Y={Y}')

copy_ps(Y, X)
print(f'point 2, X={X}, Y={Y}')

import pptx
from pptx import __version__ as v
from pptx import Presentation
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_FILL, MSO_FILL_TYPE
import os

print(f'pptx_version={v}') # 1.0.2
import pptx.shapes
import pptx.shapes.autoshape
import pptx.text
import pptx.text.text


class Prs(Presentation):
    def __init__(self, template_filepath:str, text_substitutions:list):
        if not os.path.isfile(template_filepath):
            print(f'{template_filepath} does not exists')
            exit(-1)
        self._prs = super().__init__(template_filepath)

    def _copy_font(self, src:pptx.text.text._Run, dest:pptx.text.text._Run):
        src_font = src.font
        dest_font = dest.font
        dest_font.size = src_font.size
        dest_font.bold = src_font.bold
        dest_font.italic = src_font.italic
        dest_font.underline = src_font.underline
        dest_font.color.rgb = src_font.color.rgb
        dest_font.name = src_font.name
    
    def _substitute_a_text(self, shape:pptx.shapes.autoshape.Shape, new_text:str):
        text_frame = shape.text_frame
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.text = new_text
                self._copy_font(run, run)
    
    def substitute_texts(self, text_substitutions:list):
        for shape in self.slides[0].shapes:
            if not shape.has_text_frame:
                continue
            for text_substitution in text_substitutions:
                if shape.text_frame.text == text_substitution[0]:
                    self._substitute_a_text(shape, text_substitution[1])
                    break
    
    def go_through_all_shapes(self):
        for slide in self._prs.slides:
            for shape in slide.shapes:
                print(f'{shape}')
                if shape.has_text_frame:
                    print(f'{shape.text_frame.text}')
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            print(f'{run.text}')
                            print(f'{run.font.size}')
                            print(f'{run.font.bold}')
                            print(f'{run.font.italic}')
                            print(f'{run.font.underline}')
                            print(f'{run.font.color.rgb}')
                            print(f'{run.font.name}')
    
    def save(self, output_filepath:str):
        self.save(output_filepath)
        
prs = Presentation(".\\makeMonMeetPdf\\raw_main.pptx")

for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text:
            if 'host' in shape.text:
                print(f'id={shape.shape_id},text={shape.text}')
                text_frame = shape.text_frame
                first_parag = text_frame.paragraphs[0]
                first_run = first_parag.runs[0] if first_parag.runs else first_parag.add_run()
                font = first_run.font
                font_name = font.name
                font_size = font.size
                font_bold = font.bold
                font_italic = font.italic
                font_underline = font.underline
                font_color = font.color
                text_frame.clear()
                new_run = text_frame.paragraphs[0].add_run()
                new_run.text = '曾嘉明DM'
                new_run.font.name = font_name
                new_run.font.size = font_size
                new_run.font.bold = font_bold
                new_run.font.italic = font_italic
                new_run.font.underline = font_underline
                if font_color.type == MSO_COLOR_TYPE.SCHEME:
                    print(f'font_color={font_color.theme_color}')
                    new_run.font.color.theme_color = font_color.theme_color
                else:
                    print(f'color_type={font_color.type}')
prs.save('.\\makeMonMeetPdf\\raw_main_2.pptx')
